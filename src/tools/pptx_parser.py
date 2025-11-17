"""
PowerPoint Parser for extracting content from PPTX files.

This module provides comprehensive PowerPoint parsing capabilities for:
- Text extraction from slides
- Speaker notes extraction
- Table extraction from slides
- Image detection and metadata
- Comprehensive presentation metadata
"""

import logging
from pathlib import Path
from typing import Dict, List, Any, Optional
import io

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import pandas as pd
except ImportError as e:
    raise ImportError(f"Missing required dependencies for PowerPoint parsing: {e}")

logger = logging.getLogger(__name__)


class PPTXParser:
    """
    Enterprise-grade PowerPoint presentation parser.

    Features:
    - Slide-by-slide content extraction
    - Speaker notes extraction
    - Table detection and extraction
    - Image metadata extraction
    - Shape text extraction
    - Comprehensive presentation metadata
    """

    def __init__(self, extract_images: bool = True, extract_tables: bool = True):
        """
        Initialize PowerPoint parser.

        Args:
            extract_images: Whether to extract image metadata
            extract_tables: Whether to extract tables from slides
        """
        self.extract_images = extract_images
        self.extract_tables = extract_tables

    def parse_presentation(self, pptx_path: str) -> Dict[str, Any]:
        """
        Parse entire PowerPoint presentation.

        Args:
            pptx_path: Path to PPTX file

        Returns:
            Dictionary containing:
                - metadata: Presentation metadata
                - slides: List of slide data dictionaries
                - summary: Summary statistics
        """
        try:
            pptx_path = Path(pptx_path)

            if not pptx_path.exists():
                logger.error(f"PowerPoint file not found: {pptx_path}")
                return self._empty_result(f"File not found: {pptx_path}")

            if pptx_path.suffix.lower() not in ['.pptx', '.ppt']:
                logger.error(f"Invalid file format: {pptx_path.suffix}")
                return self._empty_result(f"Invalid format: {pptx_path.suffix}")

            # Load presentation
            prs = Presentation(pptx_path)

            # Extract metadata
            metadata = self._extract_metadata(prs, pptx_path)

            # Process slides
            slides = []
            for i, slide in enumerate(prs.slides):
                slide_data = self.parse_slide(slide, i)
                slides.append(slide_data)

            # Calculate summary
            summary = self._calculate_summary(slides)

            result = {
                'metadata': metadata,
                'slides': slides,
                'summary': summary,
                'success': True,
                'error': None
            }

            logger.info(
                f"Parsed presentation: {len(slides)} slides, "
                f"{summary['total_text_length']} chars"
            )

            return result

        except Exception as e:
            logger.error(f"Error parsing PowerPoint: {e}")
            return self._empty_result(str(e))

    def parse_slide(self, slide: Any, slide_number: int) -> Dict[str, Any]:
        """
        Parse a single slide.

        Args:
            slide: PowerPoint slide object
            slide_number: Slide index (0-based)

        Returns:
            Dictionary with slide data
        """
        try:
            # Extract slide text
            text = self.extract_slide_text(slide)

            # Extract speaker notes
            notes = self.extract_slide_notes(slide)

            # Extract tables if enabled
            tables = []
            if self.extract_tables:
                tables = self.extract_slide_tables(slide)

            # Extract image metadata if enabled
            images = []
            if self.extract_images:
                images = self.extract_slide_images(slide)

            slide_data = {
                'slide_number': slide_number,
                'text': text,
                'notes': notes,
                'tables': tables,
                'images': images,
                'table_count': len(tables),
                'image_count': len(images),
                'text_length': len(text),
                'has_notes': len(notes) > 0
            }

            return slide_data

        except Exception as e:
            logger.error(f"Error parsing slide {slide_number}: {e}")
            return {
                'slide_number': slide_number,
                'error': str(e),
                'text': '',
                'notes': '',
                'tables': [],
                'images': []
            }

    def extract_slide_text(self, slide: Any) -> str:
        """
        Extract all text from a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            Concatenated text from all shapes
        """
        try:
            text_parts = []

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_parts.append(shape.text.strip())

            return "\n".join(text_parts)

        except Exception as e:
            logger.warning(f"Error extracting slide text: {e}")
            return ""

    def extract_slide_notes(self, slide: Any) -> str:
        """
        Extract speaker notes from a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            Speaker notes text
        """
        try:
            if not hasattr(slide, 'notes_slide'):
                return ""

            notes_slide = slide.notes_slide

            if not hasattr(notes_slide, 'notes_text_frame'):
                return ""

            notes_text_frame = notes_slide.notes_text_frame

            if notes_text_frame and hasattr(notes_text_frame, 'text'):
                return notes_text_frame.text.strip()

            return ""

        except Exception as e:
            logger.warning(f"Error extracting slide notes: {e}")
            return ""

    def extract_slide_images(self, slide: Any) -> List[Dict[str, Any]]:
        """
        Extract image metadata from a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            List of image metadata dictionaries
        """
        try:
            images = []

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_data = {
                        'type': 'image',
                        'name': shape.name,
                        'width': shape.width,
                        'height': shape.height,
                        'left': shape.left,
                        'top': shape.top
                    }

                    # Try to get image format
                    if hasattr(shape, 'image'):
                        image = shape.image
                        image_data['format'] = image.content_type
                        image_data['size'] = len(image.blob)

                    images.append(image_data)

            return images

        except Exception as e:
            logger.warning(f"Error extracting slide images: {e}")
            return []

    def extract_slide_tables(self, slide: Any) -> List[pd.DataFrame]:
        """
        Extract tables from a slide.

        Args:
            slide: PowerPoint slide object

        Returns:
            List of DataFrames representing tables
        """
        try:
            tables = []

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    try:
                        table_data = self._extract_table_from_shape(shape)
                        if table_data is not None and not table_data.empty:
                            tables.append(table_data)
                    except Exception as e:
                        logger.warning(f"Error extracting table from shape: {e}")
                        continue

            return tables

        except Exception as e:
            logger.warning(f"Error extracting slide tables: {e}")
            return []

    def _extract_table_from_shape(self, shape: Any) -> Optional[pd.DataFrame]:
        """
        Extract table data from a table shape.

        Args:
            shape: PowerPoint table shape

        Returns:
            DataFrame with table data
        """
        try:
            table = shape.table

            # Extract table data
            data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_text = cell.text.strip() if hasattr(cell, 'text') else ''
                    row_data.append(cell_text)
                data.append(row_data)

            if not data:
                return None

            # Create DataFrame
            df = pd.DataFrame(data)

            # Use first row as headers if it looks like a header
            if len(df) > 1:
                df.columns = df.iloc[0]
                df = df[1:]
                df = df.reset_index(drop=True)

            return df

        except Exception as e:
            logger.warning(f"Error extracting table from shape: {e}")
            return None

    def _extract_metadata(
        self,
        prs: Presentation,
        file_path: Path
    ) -> Dict[str, Any]:
        """
        Extract presentation metadata.

        Args:
            prs: Presentation object
            file_path: Path to presentation file

        Returns:
            Dictionary with metadata
        """
        try:
            core_props = prs.core_properties

            metadata = {
                'filename': file_path.name,
                'filepath': str(file_path),
                'slide_count': len(prs.slides),
                'slide_width': prs.slide_width,
                'slide_height': prs.slide_height,
                'title': core_props.title or '',
                'author': core_props.author or '',
                'subject': core_props.subject or '',
                'keywords': core_props.keywords or '',
                'comments': core_props.comments or '',
                'created': str(core_props.created) if core_props.created else None,
                'modified': str(core_props.modified) if core_props.modified else None,
                'last_modified_by': core_props.last_modified_by or ''
            }

            return metadata

        except Exception as e:
            logger.warning(f"Error extracting metadata: {e}")
            return {
                'filename': file_path.name,
                'filepath': str(file_path)
            }

    def _calculate_summary(self, slides: List[Dict[str, Any]]) -> Dict[str, Any]:
        """
        Calculate summary statistics from slides.

        Args:
            slides: List of slide data dictionaries

        Returns:
            Dictionary with summary statistics
        """
        try:
            total_text_length = sum(s.get('text_length', 0) for s in slides)
            total_tables = sum(s.get('table_count', 0) for s in slides)
            total_images = sum(s.get('image_count', 0) for s in slides)
            slides_with_notes = sum(1 for s in slides if s.get('has_notes', False))

            # Concatenate all text
            all_text = "\n\n".join(s.get('text', '') for s in slides)

            # Concatenate all notes
            all_notes = "\n\n".join(s.get('notes', '') for s in slides if s.get('notes'))

            return {
                'total_slides': len(slides),
                'total_text_length': total_text_length,
                'total_tables': total_tables,
                'total_images': total_images,
                'slides_with_notes': slides_with_notes,
                'avg_text_per_slide': total_text_length / len(slides) if slides else 0,
                'all_text': all_text,
                'all_notes': all_notes
            }

        except Exception as e:
            logger.error(f"Error calculating summary: {e}")
            return {
                'total_slides': len(slides),
                'error': str(e)
            }

    def extract_all_text(self, pptx_path: str) -> str:
        """
        Extract all text from presentation (convenience method).

        Args:
            pptx_path: Path to PPTX file

        Returns:
            All text concatenated
        """
        try:
            result = self.parse_presentation(pptx_path)

            if not result['success']:
                return ""

            summary = result.get('summary', {})
            all_text = summary.get('all_text', '')
            all_notes = summary.get('all_notes', '')

            # Combine slide text and notes
            combined = f"{all_text}\n\n--- Notes ---\n\n{all_notes}" if all_notes else all_text

            return combined

        except Exception as e:
            logger.error(f"Error extracting all text: {e}")
            return ""

    def extract_all_tables(self, pptx_path: str) -> List[pd.DataFrame]:
        """
        Extract all tables from presentation (convenience method).

        Args:
            pptx_path: Path to PPTX file

        Returns:
            List of all tables as DataFrames
        """
        try:
            result = self.parse_presentation(pptx_path)

            if not result['success']:
                return []

            all_tables = []
            for slide in result.get('slides', []):
                tables = slide.get('tables', [])
                all_tables.extend(tables)

            return all_tables

        except Exception as e:
            logger.error(f"Error extracting all tables: {e}")
            return []

    def get_slide_count(self, pptx_path: str) -> int:
        """
        Get number of slides in presentation.

        Args:
            pptx_path: Path to PPTX file

        Returns:
            Number of slides
        """
        try:
            prs = Presentation(pptx_path)
            return len(prs.slides)

        except Exception as e:
            logger.error(f"Error getting slide count: {e}")
            return 0

    def _empty_result(self, error: str) -> Dict[str, Any]:
        """
        Return empty result with error message.

        Args:
            error: Error message

        Returns:
            Empty result dictionary
        """
        return {
            'metadata': {},
            'slides': [],
            'summary': {'total_slides': 0},
            'success': False,
            'error': error
        }
