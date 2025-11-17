"""
Comprehensive tests for PowerPoint Parser.

Tests cover:
- Presentation parsing
- Slide text extraction
- Speaker notes extraction
- Table extraction
- Image metadata extraction
- Metadata extraction
- Error handling
"""

import pytest
from unittest.mock import Mock, patch, MagicMock, PropertyMock
from pathlib import Path
import pandas as pd

from src.tools.pptx_parser import PPTXParser


@pytest.fixture
def pptx_parser():
    """Create PowerPoint parser instance."""
    return PPTXParser(extract_images=True, extract_tables=True)


@pytest.fixture
def mock_presentation():
    """Create a mock Presentation object."""
    prs = MagicMock()
    prs.slides = []
    prs.slide_width = 10000000
    prs.slide_height = 7500000

    # Mock core properties
    core_props = Mock()
    core_props.title = "Test Presentation"
    core_props.author = "Test Author"
    core_props.subject = "Test Subject"
    core_props.keywords = "test, keywords"
    core_props.comments = "Test comments"
    core_props.created = None
    core_props.modified = None
    core_props.last_modified_by = "Test User"

    prs.core_properties = core_props

    return prs


@pytest.fixture
def mock_slide():
    """Create a mock slide object."""
    slide = MagicMock()
    slide.shapes = []

    # Mock notes slide
    notes_slide = Mock()
    notes_text_frame = Mock()
    notes_text_frame.text = "Test speaker notes"
    notes_slide.notes_text_frame = notes_text_frame
    slide.notes_slide = notes_slide

    return slide


@pytest.fixture
def mock_text_shape():
    """Create a mock text shape."""
    shape = Mock()
    shape.text = "Test slide text"
    shape.has_text_frame = True
    return shape


class TestPPTXParserInit:
    """Test PowerPoint parser initialization."""

    def test_init_default_params(self):
        """Test initialization with default parameters."""
        parser = PPTXParser()
        assert parser.extract_images is True
        assert parser.extract_tables is True

    def test_init_custom_params(self):
        """Test initialization with custom parameters."""
        parser = PPTXParser(extract_images=False, extract_tables=False)
        assert parser.extract_images is False
        assert parser.extract_tables is False


class TestParsePresentation:
    """Test full presentation parsing."""

    @patch('src.tools.pptx_parser.Presentation')
    def test_parse_presentation_success(
        self,
        mock_prs_class,
        pptx_parser,
        mock_presentation,
        mock_slide,
        tmp_path
    ):
        """Test successful presentation parsing."""
        test_file = tmp_path / "test.pptx"
        test_file.touch()

        mock_presentation.slides = [mock_slide]
        mock_prs_class.return_value = mock_presentation

        result = pptx_parser.parse_presentation(str(test_file))

        assert result['success'] is True
        assert 'metadata' in result
        assert 'slides' in result
        assert 'summary' in result
        assert len(result['slides']) == 1

    def test_parse_presentation_file_not_found(self, pptx_parser):
        """Test parsing with non-existent file."""
        result = pptx_parser.parse_presentation("/nonexistent/file.pptx")

        assert result['success'] is False
        assert 'not found' in result['error'].lower()

    @patch('src.tools.pptx_parser.Presentation')
    def test_parse_presentation_invalid_format(self, mock_prs_class, pptx_parser, tmp_path):
        """Test parsing with invalid file format."""
        test_file = tmp_path / "test.xyz"
        test_file.touch()

        result = pptx_parser.parse_presentation(str(test_file))

        assert result['success'] is False
        assert 'format' in result['error'].lower()

    @patch('src.tools.pptx_parser.Presentation')
    def test_parse_presentation_exception_handling(
        self,
        mock_prs_class,
        pptx_parser,
        tmp_path
    ):
        """Test exception handling during parsing."""
        test_file = tmp_path / "test.pptx"
        test_file.touch()

        mock_prs_class.side_effect = Exception("Parse error")

        result = pptx_parser.parse_presentation(str(test_file))

        assert result['success'] is False
        assert result['error'] == "Parse error"


class TestParseSlide:
    """Test individual slide parsing."""

    def test_parse_slide_with_text(self, pptx_parser, mock_slide, mock_text_shape):
        """Test parsing slide with text."""
        mock_slide.shapes = [mock_text_shape]

        slide_data = pptx_parser.parse_slide(mock_slide, 0)

        assert slide_data['slide_number'] == 0
        assert 'text' in slide_data
        assert 'notes' in slide_data
        assert isinstance(slide_data['tables'], list)
        assert isinstance(slide_data['images'], list)

    def test_parse_slide_exception_handling(self, pptx_parser):
        """Test exception handling during slide parsing."""
        bad_slide = Mock()
        bad_slide.shapes = None  # Will cause error

        slide_data = pptx_parser.parse_slide(bad_slide, 0)

        assert 'error' in slide_data
        assert slide_data['text'] == ''


class TestExtractSlideText:
    """Test slide text extraction."""

    def test_extract_text_single_shape(self, pptx_parser, mock_slide, mock_text_shape):
        """Test extraction with single text shape."""
        mock_slide.shapes = [mock_text_shape]

        text = pptx_parser.extract_slide_text(mock_slide)

        assert "Test slide text" in text

    def test_extract_text_multiple_shapes(self, pptx_parser, mock_slide):
        """Test extraction with multiple text shapes."""
        shape1 = Mock()
        shape1.text = "First text"
        shape2 = Mock()
        shape2.text = "Second text"

        mock_slide.shapes = [shape1, shape2]

        text = pptx_parser.extract_slide_text(mock_slide)

        assert "First text" in text
        assert "Second text" in text

    def test_extract_text_empty_slide(self, pptx_parser, mock_slide):
        """Test extraction from empty slide."""
        mock_slide.shapes = []

        text = pptx_parser.extract_slide_text(mock_slide)

        assert text == ""

    def test_extract_text_shape_without_text(self, pptx_parser, mock_slide):
        """Test extraction from shape without text attribute."""
        shape = Mock(spec=[])  # No 'text' attribute

        mock_slide.shapes = [shape]

        text = pptx_parser.extract_slide_text(mock_slide)

        assert text == ""


class TestExtractSlideNotes:
    """Test speaker notes extraction."""

    def test_extract_notes_success(self, pptx_parser, mock_slide):
        """Test successful notes extraction."""
        notes = pptx_parser.extract_slide_notes(mock_slide)

        assert notes == "Test speaker notes"

    def test_extract_notes_no_notes_slide(self, pptx_parser):
        """Test extraction with no notes slide."""
        slide = Mock(spec=[])  # No notes_slide attribute

        notes = pptx_parser.extract_slide_notes(slide)

        assert notes == ""

    def test_extract_notes_no_text_frame(self, pptx_parser):
        """Test extraction with no notes text frame."""
        slide = Mock()
        slide.notes_slide = Mock(spec=[])  # No notes_text_frame

        notes = pptx_parser.extract_slide_notes(slide)

        assert notes == ""


class TestExtractSlideImages:
    """Test image metadata extraction."""

    def test_extract_images_with_picture(self, pptx_parser, mock_slide):
        """Test extraction with picture shapes."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        picture_shape = Mock()
        picture_shape.shape_type = MSO_SHAPE_TYPE.PICTURE
        picture_shape.name = "Picture 1"
        picture_shape.width = 1000
        picture_shape.height = 800
        picture_shape.left = 100
        picture_shape.top = 100

        # Mock image property
        mock_image = Mock()
        mock_image.content_type = "image/png"
        mock_image.blob = b"fake_image_data"
        picture_shape.image = mock_image

        mock_slide.shapes = [picture_shape]

        images = pptx_parser.extract_slide_images(mock_slide)

        assert len(images) == 1
        assert images[0]['type'] == 'image'
        assert images[0]['name'] == "Picture 1"
        assert images[0]['format'] == "image/png"

    def test_extract_images_no_pictures(self, pptx_parser, mock_slide):
        """Test extraction with no pictures."""
        mock_slide.shapes = []

        images = pptx_parser.extract_slide_images(mock_slide)

        assert images == []

    def test_extract_images_disabled(self, mock_slide):
        """Test extraction when images disabled."""
        parser = PPTXParser(extract_images=False)
        # Still should return empty list when method called directly
        images = parser.extract_slide_images(mock_slide)

        assert isinstance(images, list)


class TestExtractSlideTables:
    """Test table extraction from slides."""

    def test_extract_tables_with_table_shape(self, pptx_parser, mock_slide):
        """Test extraction with table shapes."""
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        table_shape = Mock()
        table_shape.shape_type = MSO_SHAPE_TYPE.TABLE

        # Mock table with cells
        mock_table = Mock()
        mock_row1 = Mock()
        mock_row2 = Mock()

        cell1 = Mock()
        cell1.text = "Header"
        cell2 = Mock()
        cell2.text = "Data"

        mock_row1.cells = [cell1]
        mock_row2.cells = [cell2]

        mock_table.rows = [mock_row1, mock_row2]
        table_shape.table = mock_table

        mock_slide.shapes = [table_shape]

        tables = pptx_parser.extract_slide_tables(mock_slide)

        assert len(tables) > 0
        assert isinstance(tables[0], pd.DataFrame)

    def test_extract_tables_no_tables(self, pptx_parser, mock_slide):
        """Test extraction with no tables."""
        mock_slide.shapes = []

        tables = pptx_parser.extract_slide_tables(mock_slide)

        assert tables == []

    def test_extract_tables_disabled(self, mock_slide):
        """Test extraction when tables disabled."""
        parser = PPTXParser(extract_tables=False)
        # Still should return empty list when method called directly
        tables = parser.extract_slide_tables(mock_slide)

        assert isinstance(tables, list)


class TestExtractMetadata:
    """Test metadata extraction."""

    def test_extract_metadata_success(self, pptx_parser, mock_presentation, tmp_path):
        """Test successful metadata extraction."""
        test_file = tmp_path / "test.pptx"

        metadata = pptx_parser._extract_metadata(mock_presentation, test_file)

        assert metadata['filename'] == "test.pptx"
        assert metadata['title'] == "Test Presentation"
        assert metadata['author'] == "Test Author"
        assert metadata['slide_count'] == 0  # No slides in mock

    def test_extract_metadata_exception_handling(self, pptx_parser, tmp_path):
        """Test exception handling in metadata extraction."""
        bad_prs = Mock()
        bad_prs.core_properties = None  # Will cause error
        bad_prs.slides = []

        test_file = tmp_path / "test.pptx"

        metadata = pptx_parser._extract_metadata(bad_prs, test_file)

        # Should still return basic metadata
        assert 'filename' in metadata
        assert 'filepath' in metadata


class TestCalculateSummary:
    """Test summary calculation."""

    def test_calculate_summary_normal(self, pptx_parser):
        """Test summary calculation with normal data."""
        slides = [
            {
                'text': 'Slide 1 text',
                'notes': 'Notes 1',
                'text_length': 12,
                'table_count': 1,
                'image_count': 2,
                'has_notes': True
            },
            {
                'text': 'Slide 2 text',
                'notes': '',
                'text_length': 12,
                'table_count': 0,
                'image_count': 1,
                'has_notes': False
            }
        ]

        summary = pptx_parser._calculate_summary(slides)

        assert summary['total_slides'] == 2
        assert summary['total_text_length'] == 24
        assert summary['total_tables'] == 1
        assert summary['total_images'] == 3
        assert summary['slides_with_notes'] == 1
        assert summary['avg_text_per_slide'] == 12

    def test_calculate_summary_empty(self, pptx_parser):
        """Test summary calculation with no slides."""
        summary = pptx_parser._calculate_summary([])

        assert summary['total_slides'] == 0
        assert summary['avg_text_per_slide'] == 0


class TestConvenienceMethods:
    """Test convenience methods."""

    @patch('src.tools.pptx_parser.Presentation')
    def test_extract_all_text(
        self,
        mock_prs_class,
        pptx_parser,
        mock_presentation,
        mock_slide,
        mock_text_shape,
        tmp_path
    ):
        """Test extracting all text from presentation."""
        test_file = tmp_path / "test.pptx"
        test_file.touch()

        mock_slide.shapes = [mock_text_shape]
        mock_presentation.slides = [mock_slide]
        mock_prs_class.return_value = mock_presentation

        text = pptx_parser.extract_all_text(str(test_file))

        assert isinstance(text, str)
        assert len(text) > 0

    @patch('src.tools.pptx_parser.Presentation')
    def test_extract_all_tables(
        self,
        mock_prs_class,
        pptx_parser,
        mock_presentation,
        tmp_path
    ):
        """Test extracting all tables from presentation."""
        test_file = tmp_path / "test.pptx"
        test_file.touch()

        mock_presentation.slides = []
        mock_prs_class.return_value = mock_presentation

        tables = pptx_parser.extract_all_tables(str(test_file))

        assert isinstance(tables, list)

    @patch('src.tools.pptx_parser.Presentation')
    def test_get_slide_count(
        self,
        mock_prs_class,
        pptx_parser,
        mock_presentation,
        tmp_path
    ):
        """Test getting slide count."""
        test_file = tmp_path / "test.pptx"
        test_file.touch()

        mock_presentation.slides = [Mock(), Mock(), Mock()]
        mock_prs_class.return_value = mock_presentation

        count = pptx_parser.get_slide_count(str(test_file))

        assert count == 3

    def test_get_slide_count_exception(self, pptx_parser):
        """Test getting slide count with exception."""
        count = pptx_parser.get_slide_count("/nonexistent/file.pptx")

        assert count == 0


class TestExtractTableFromShape:
    """Test table extraction from shape."""

    def test_extract_table_success(self, pptx_parser):
        """Test successful table extraction."""
        shape = Mock()
        mock_table = Mock()

        # Create mock table with rows
        row1 = Mock()
        row2 = Mock()

        cell1 = Mock()
        cell1.text = "Name"
        cell2 = Mock()
        cell2.text = "Alice"

        row1.cells = [cell1]
        row2.cells = [cell2]

        mock_table.rows = [row1, row2]
        shape.table = mock_table

        df = pptx_parser._extract_table_from_shape(shape)

        assert df is not None
        assert isinstance(df, pd.DataFrame)

    def test_extract_table_empty(self, pptx_parser):
        """Test extraction from empty table."""
        shape = Mock()
        mock_table = Mock()
        mock_table.rows = []
        shape.table = mock_table

        df = pptx_parser._extract_table_from_shape(shape)

        assert df is None

    def test_extract_table_exception(self, pptx_parser):
        """Test exception handling in table extraction."""
        shape = Mock()
        shape.table = None  # Will cause error

        df = pptx_parser._extract_table_from_shape(shape)

        assert df is None


class TestEmptyResult:
    """Test empty result generation."""

    def test_empty_result(self, pptx_parser):
        """Test empty result creation."""
        result = pptx_parser._empty_result("Test error")

        assert result['success'] is False
        assert result['error'] == "Test error"
        assert result['metadata'] == {}
        assert result['slides'] == []
