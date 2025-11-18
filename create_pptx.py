from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Cybersecurity Risk Assessment"
subtitle.text = "Q4 2024 Executive Summary"

# Slide 2: Risk Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Key Risk Findings"
tf = content.text_frame
tf.text = "Critical Vulnerabilities Identified"
p = tf.add_paragraph()
p.text = "5 High-severity CVEs require immediate remediation"
p.level = 1
p = tf.add_paragraph()
p.text = "Affected systems: Web servers, database clusters"
p.level = 1
p = tf.add_paragraph()
p.text = "MITRE ATT&CK: T1190 (Exploit Public-Facing Application)"
p.level = 1

# Slide 3: Risk Matrix
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Risk Scoring Matrix"

# Add table
rows, cols = 4, 3
left = Inches(2)
top = Inches(2.5)
width = Inches(6)
height = Inches(3)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Headers
table.cell(0, 0).text = "Risk ID"
table.cell(0, 1).text = "Likelihood"
table.cell(0, 2).text = "Impact"

# Data
table.cell(1, 0).text = "CVE-2024-1234"
table.cell(1, 1).text = "High"
table.cell(1, 2).text = "Critical"

table.cell(2, 0).text = "CVE-2024-5678"
table.cell(2, 1).text = "Medium"
table.cell(2, 2).text = "High"

table.cell(3, 0).text = "CVE-2024-9012"
table.cell(3, 1).text = "Low"
table.cell(3, 2).text = "Medium"

# Slide 4: Recommendations
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Remediation Recommendations"
tf = content.text_frame
tf.text = "Immediate Actions"
p = tf.add_paragraph()
p.text = "Patch web application servers within 48 hours"
p.level = 1
p = tf.add_paragraph()
p.text = "Implement WAF rules for CVE-2024-1234"
p.level = 1
p = tf.add_paragraph()
p.text = "Conduct vulnerability re-scan post-remediation"
p.level = 1

# Save
prs.save('sample_security_presentation.pptx')
print("Created: sample_security_presentation.pptx")
